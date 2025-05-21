# File: apps/pkm-indexer/organize.py
import os
import shutil
import time
import frontmatter
import openai
import re
import json
from pathlib import Path
import pdfplumber
import pytesseract
from PIL import Image
import requests
from bs4 import BeautifulSoup
import logging # Add logging import if not already there

openai.api_key = os.getenv("OPENAI_API_KEY")
logger = logging.getLogger(__name__) # Or use your existing logger

def infer_file_type(filename):
    ext = Path(filename).suffix.lower().strip() # Added .strip()
    if ext in [".md", ".txt"]: return "text"
    if ext in [".pdf"]: return "pdf"
    if ext in [".png", ".jpg", ".jpeg", ".gif", ".bmp"]: return "image"
    if ext in [".mp3", ".wav", ".m4a"]: return "audio"
    if ext in [".doc", ".docx"]: return "document"
    if ext in [".ppt", ".pptx"]: return "presentation"
    if ext in [".xls", ".xlsx"]: return "spreadsheet"
    if ext in [".rtf"]: return "rtf"
    return "other"

def extract_text_from_pdf(path):
    try:
        with pdfplumber.open(path) as pdf:
            text = "\n".join(page.extract_text() or "" for page in pdf.pages)
            
            # Check if this looks like a LinkedIn post
            if "Profile viewers" in text[:500] or "Post impressions" in text[:500] or "linkedin.com" in text.lower():
                return process_linkedin_pdf(text, path)
            
            return text
    except Exception as e:
        return f"[PDF extraction failed: {e}]"
    
def extract_text_from_docx(path):
    """Extract text content from a .docx file."""
    try:
        import docx
        doc = docx.Document(path)
        full_text = []
        
        # Extract text from paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                full_text.append(para.text)
        
        # Extract text from tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        full_text.append(cell.text)
        
        return "\n\n".join(full_text)
    except Exception as e:
        return f"[Word document extraction failed: {e}]"

def extract_text_from_pptx(path):
    """Extract text content from a .pptx file."""
    try:
        import pptx
        prs = pptx.Presentation(path)
        full_text = []
        
        for i, slide in enumerate(prs.slides):
            slide_text = [f"--- Slide {i+1} ---"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            
            if len(slide_text) > 1:  # Only add slides with actual content
                full_text.append("\n\n".join(slide_text))
        
        return "\n\n".join(full_text)
    except Exception as e:
        return f"[PowerPoint extraction failed: {e}]"

def extract_text_from_xlsx(path):
    """Extract text content from an .xlsx file."""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        full_text = []
        
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            sheet_text = [f"--- Sheet: {sheet_name} ---"]
            
            for row in sheet.iter_rows(values_only=True):
                row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                if row_text.strip():
                    sheet_text.append(row_text)
            
            if len(sheet_text) > 1:  # Only add sheets with actual content
                full_text.append("\n\n".join(sheet_text))
        
        return "\n\n".join(full_text)
    except Exception as e:
        return f"[Excel extraction failed: {e}]"

def extract_text_from_markdown(path):
    """Extract text content from .md files, preserving the original markdown."""
    try:
        with open(path, 'r', encoding='utf-8') as file:
            content = file.read()
        # Just return the raw markdown - it's already text
        return content
    except Exception as e:
        return f"[Markdown extraction failed: {e}]"

def extract_text_from_rtf(path):
    """Extract text content from .rtf files."""
    try:
        with open(path, 'r', encoding='utf-8', errors='ignore') as file:
            content = file.read()
            
        # Basic RTF parsing to remove control codes
        # For better parsing, we'd use the python-rtf library
        # This is a simple fallback in case the library fails
        text = re.sub(r'\\[a-z]+(-?\d+)?[ ]?', ' ', content)
        text = re.sub(r'\{|\}|\\|\|', '', text)
        
        return text
    except Exception as e:
        try:
            # If basic parsing fails, try using the python-rtf library
            import striprtf
            with open(path, 'r', encoding='utf-8', errors='ignore') as file:
                content = file.read()
            return striprtf.rtf_to_text(content)
        except Exception as rtf_error:
            return f"[RTF extraction failed: {e}, {rtf_error}]"

def process_linkedin_pdf(text, path):
    """Process LinkedIn PDF content to extract the main post and ignore comments."""
    try:
        # Pattern to detect the start of comments section
        comment_indicators = [
            "Reactions", 
            "Like · Reply",
            "comments · ",
            "reposts",
            "Most relevant"
        ]
        
        # Pattern to detect URLs in the text
        url_pattern = r'https?://(?:[-\w.]|(?:%[\da-fA-F]{2}))+(?:/[-\w%!.~\'()*+,;=:@/&?=]*)?'
        
        # Find all URLs in the original content
        main_urls = re.findall(url_pattern, text)
        important_urls = []
        
        # Split content by lines to process
        lines = text.split('\n')
        main_content_lines = []
        
        # Track if we're in the comments section
        in_comments = False
        author_comment_section = False
        post_author = None
        
        # Extract the post author if available (usually near the beginning)
        for i, line in enumerate(lines[:15]):
            if "• Author" in line or "• 1st" in line or "• 2nd" in line or "• 3rd" in line:
                # The line before often contains the author name
                if i > 0:
                    post_author = lines[i-1].strip()
                    break
                    
        # Process the content line by line
        for i, line in enumerate(lines):
            # Check if we've hit the comments section
            if any(indicator in line for indicator in comment_indicators) and i > 10:
                in_comments = True
                continue
                
            # If we're still in the main content, keep the line
            if not in_comments:
                main_content_lines.append(line)
                continue
                
            # Check for author comments (only process if we know the author)
            if post_author and post_author in line and i+2 < len(lines) and "Author" in lines[i:i+2]:
                author_comment_section = True
                continue
                
            # Process author comment content
            if author_comment_section:
                # Look for URLs or other important info in first author comment
                urls_in_comment = re.findall(url_pattern, line)
                if urls_in_comment:
                    important_urls.extend(urls_in_comment)
                    
                # Check if author comment section is ending
                if "Like · Reply" in line or "Like · " in line:
                    author_comment_section = False
        
        # Combine the main content
        main_content = '\n'.join(main_content_lines)
        
        # Add any important URLs from author comments if they weren't in the main content
        for url in important_urls:
            if url not in main_urls and ("lnkd.in" in url or ".com" in url):  # LinkedIn short URLs are often important
                main_content += f"\n\nAdditional URL from author comment: {url}"
                
        print("📱 Detected LinkedIn post, removed comments section")
        return main_content
    except Exception as e:
        print(f"Error processing LinkedIn content: {e}")
        return text  # Return original if processing fails

def extract_text_from_image(path):
    try:
        # Open and process image
        image = Image.open(path)
        
        # Try multiple preprocessing approaches
        texts = []
        
        # Approach 1: Original with adjusted threshold
        img1 = image.convert("L")
        img1 = img1.point(lambda x: 0 if x < 120 else 255)  # Lowered threshold
        texts.append(pytesseract.image_to_string(img1, lang="eng"))
        
        # Approach 2: Try Danish language if available
        try:
            texts.append(pytesseract.image_to_string(img1, lang="dan"))
        except:
            # If Danish not installed, try with English
            pass
        
        # Approach 3: Try with different preprocessing
        img3 = image.convert("L")
        img3 = img3.resize((int(img3.width * 1.5), int(img3.height * 1.5)), Image.LANCZOS)  # Upsample
        img3 = img3.point(lambda x: 0 if x < 150 else 255)  # Different threshold
        
        # Try multilingual if available
        try:
            texts.append(pytesseract.image_to_string(img3, lang="dan+eng"))
        except:
            texts.append(pytesseract.image_to_string(img3, lang="eng"))
        
        # Approach 4: Higher contrast for slide presentations
        img4 = image.convert("L")
        # Apply more aggressive contrast for presentation slides
        img4 = img4.point(lambda x: 0 if x < 180 else 255)
        texts.append(pytesseract.image_to_string(img4, lang="eng"))
        
        # Use the longest text result that isn't just garbage
        valid_texts = [t for t in texts if len(t.strip()) > 20]
        if valid_texts:
            text = max(valid_texts, key=len)
        else:
            text = max(texts, key=len)
        
        # If we got nothing meaningful, report failure
        if len(text.strip()) < 20:
            return "[OCR produced insufficient text. Manual processing recommended.]"
            
        print("🖼️ OCR output:", repr(text[:500]))
        return text
    except Exception as e:
        return f"[OCR failed: {e}]"

def extract_urls(text):
    """
    Extract URLs from text, including both standard http/https URLs and potential 
    title-based references that might be links.
    """
    # Standard URL pattern
    url_pattern = r'https?://(?:[-\w.]|(?:%[\da-fA-F]{2}))+(?:/[-\w%!.~\'()*+,;=:@/&?=]*)?'
    urls = re.findall(url_pattern, text)
    
    # Also look for linked text with URLs like [text](url)
    markdown_links = re.findall(r'\[([^\]]+)\]\(([^)]+)\)', text)
    markdown_urls = [link[1] for link in markdown_links if link[1].startswith('http')]
    
    # Look for potential title links in specific formats
    potential_links = []
    
    # Look for titles that might be links (for PDF resources lists)
    # Pattern: title followed by "by Author" - common in resource lists
    title_pattern = r'(?:^|\n)(?:\d+\)|\-)\s*([^""\n]+?)(?= by | \()'
    potential_titles = re.findall(title_pattern, text)
    
    # Also look for text that appears to be a clickable reference
    # Common in PDFs with links that don't have explicit URLs
    reference_patterns = [
        r'([A-Z][a-z]+(?:\s[A-Z][a-z]+)*\s(?:AI|ML|for\sEveryone|Intelligence|Awareness|Machine|clone))',  # AI-related titles
        r'(my book|my AI clone|Appendix [A-Z]|Foundry from HBS)',  # References to books, appendices, etc.
        r'(?<=see\s)([^\.,:;\n]+)'  # Things after "see" are often references
    ]
    
    for pattern in reference_patterns:
        found = re.findall(pattern, text)
        potential_links.extend([link.strip() for link in found if len(link.strip()) > 5])
    
    # Add potential titles that look like resources
    potential_links.extend([title.strip() for title in potential_titles if len(title.strip()) > 5])
    
    # Remove duplicates and very common words that aren't likely to be meaningful links
    potential_links = list(set(potential_links))
    filtered_links = [link for link in potential_links if link.lower() not in 
                     ['and', 'the', 'this', 'that', 'with', 'from', 'after', 'before']]
    
    all_urls = list(set(urls + markdown_urls))  # Remove duplicates
    print("🔗 URLs detected:", all_urls)
    print("🔍 Potential link titles:", filtered_links[:15])
    
    return all_urls, filtered_links

def enrich_urls(urls, potential_titles=None):
    enriched = []
    metadata = {}
    
    # Create a mapping of potential titles to improve URL descriptions
    title_map = {}
    if potential_titles:
        for title in potential_titles:
            # Store lowercase version for case-insensitive matching
            title_map[title.lower()] = title
    
    for url in urls:
        try:
            headers = {
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
            }
            r = requests.get(url, timeout=10, headers=headers)
            soup = BeautifulSoup(r.text, "html.parser")
            
            # Try to get title
            title = "(No title)"
            if soup.title and soup.title.string:
                title = soup.title.string.strip()
            
            # Try to get description
            description = ""
            meta_desc = soup.find('meta', attrs={'name': 'description'}) or soup.find('meta', attrs={'property': 'og:description'})
            if meta_desc and 'content' in meta_desc.attrs:
                description = meta_desc['content'].strip()
                if len(description) > 150:
                    description = description[:150] + "..."
            
            # Check if this URL might match a potential title we found
            url_lower = url.lower()
            matching_title = None
            for potential_title_lower, original_title in title_map.items():
                # See if any words from the potential title appear in the URL
                words = potential_title_lower.split()
                if any(word in url_lower for word in words if len(word) > 3):
                    matching_title = original_title
                    break
            
            # Use matching title if found
            if matching_title and len(matching_title) > 5:
                display_title = matching_title
            else:
                display_title = title
                
            enriched_entry = f"- [{display_title}]({url})"
            if description:
                enriched_entry += f"\n  *{description}*"
                
            enriched.append(enriched_entry)
            
            # Store metadata for later use
            metadata[url] = {
                "title": display_title,
                "description": description[:150] if description else "",
                "url": url
            }
            
        except Exception as e:
            enriched.append(f"- {url} (unreachable: {str(e)[:50]})")
            metadata[url] = {
                "title": url,
                "description": f"Error: {str(e)[:50]}",
                "url": url
            }
    
    print("🔍 Enriched URLs block:\\n", "\\n".join(enriched))
    return "\\n".join(enriched), metadata

def get_extract(content, file_type=None, urls_metadata=None, log_f=None, is_linkedin=False, source_filename_for_logging="UnknownFile"): # Added source_filename_for_logging
    logger.info(f"Enter get_extract for '{source_filename_for_logging}', file_type: {file_type}, is_linkedin: {is_linkedin}")
    try:
        # Check if OpenAI API key is configured
        api_key = os.getenv("OPENAI_API_KEY")
        if not api_key or not api_key.strip():
            error_msg = "OpenAI API key not configured. Please set OPENAI_API_KEY environment variable."
            logger.error(f"OpenAI API key is missing or invalid for '{source_filename_for_logging}'. Key: '{api_key}'") # MODIFIED
            # Removed log_f usage
            return "Missing API Key", "Extract failed: OpenAI API key not configured. Please add OPENAI_API_KEY to environment variables.", ["extraction_failed"]
        
        # Make sure the API key is set in openai module
        openai.api_key = api_key
        
        logger.debug(f"Content sent to GPT for '{source_filename_for_logging}' (preview): {content[:500]}") # MODIFIED from print
        
        # Determine appropriate extract length based on content
        content_length = len(content)
        if content_length < 1000:
            # For very short content, keep extract concise
            extract_length = 200
            model = "gpt-4"
        elif content_length < 5000:
            # For medium content, medium extract
            extract_length = 500
            model = "gpt-4"
        else:
            # For longer, complex content, allow longer extracts
            extract_length = 2000  # Up to ~400 words for complex content
            model = "gpt-4"  # Better for complex content
        
        # Check if this appears to be a resource list/links-heavy doc
        has_resource_patterns = (
            "resources" in content.lower() and 
            (content.count("\n1)") > 1 or content.count("\n2)") > 1)
        )
        
        # Different prompt based on content type
        if has_resource_patterns:
            # For resource-list style documents
            prompt = (
                "You are analyzing a document that appears to be a resource list with references, links, and learning materials.\n\n"
                "Create a detailed summary that specifically includes ALL referenced resources, people, and links. "
                "Also provide relevant tags that capture the subject matter and type of resources.\n\n"
                "In your extract, make sure to preserve:\n"
                "1. All resource names and titles\n"
                "2. All author names and affiliations\n"
                "3. All categories of resources\n"
                "4. Any referenced websites, tools, or platforms\n\n"
                "Respond in this JSON format:\n"
                "{\n  \"extract_title\": \"...\",\n  \"extract_content\": \"...\",\n  \"tags\": [\"tag1\", \"tag2\"]\n}\n\n"
                f"Content:\n{content[:5000]}"
            )
        elif is_linkedin:
            prompt = (
                "You are analyzing a LinkedIn post. Create a clear title and detailed summary that captures "
                "the key points, insights, and any URLs/resources mentioned in the post. Ignore promotional content.\n\n"
                "Focus on what makes this post valuable for knowledge management purposes.\n\n"
                "Respond in this JSON format:\n"
                "{\n  \"extract_title\": \"...\",\n  \"extract_content\": \"...\",\n  \"tags\": [\"tag1\", \"tag2\"]\n}\n\n"
                f"LinkedIn Post Content:\n{content[:5000]}"
            )
        elif file_type == "image":
            prompt = (
                "You are analyzing text extracted from an image via OCR. The text may have errors or be incomplete.\n\n"
                "Create a meaningful title and summary of what this image contains, plus relevant tags.\n\n"
                "For complex content, provide a detailed summary that captures the key information.\n\n"
                "Respond in this JSON format:\n"
                "{\n  \"extract_title\": \"...\",\n  \"extract_content\": \"...\",\n  \"tags\": [\"tag1\", \"tag2\"]\n}\n\n"
                f"OCR Text:\n{content[:5000]}"
            )
        elif urls_metadata and len(urls_metadata) > 0:
            # Create a summary of URLs for the prompt
            url_summary = "\n".join([f"- {data['title']}: {data['url']}" for url, data in urls_metadata.items()])
            
            prompt = (
                "You are summarizing content that contains valuable URLs and references.\n\n"
                "Create a title and detailed summary preserving key information, plus relevant tags.\n"
                "For rich content with many references, provide a comprehensive summary.\n\n"
                "Pay special attention to these detected URLs and resources:\n\n"
                f"{url_summary}\n\n"
                "Respond in this JSON format:\n"
                "{\n  \"extract_title\": \"...\",\n  \"extract_content\": \"...\",\n  \"tags\": [\"tag1\", \"tag2\"]\n}\n\n"
                f"Content:\n{content[:5000]}"
            )
        else:
            prompt = (
                "You are a semantic summarizer. Return a short title and a deeper thematic summary, plus relevant tags.\n\n"
                "For complex or information-rich content, provide a detailed summary that captures the key points.\n\n"
                "Respond in this JSON format:\n"
                "{\n  \"extract_title\": \"...\",\n  \"extract_content\": \"...\",\n  \"tags\": [\"tag1\", \"tag2\"]\n}\n\n"
                f"Content:\n{content[:5000]}"
            )
        
        # Log the prompt for debugging
        # MODIFIED: Replaced log_f with logger
        logger.debug(f"OpenAI Prompt for '{source_filename_for_logging}' (preview): {prompt[:500]}...")
        
        # Add retries for API call reliability
        max_retries = 3
        retry_delay = 2  # seconds
        
        for attempt in range(max_retries):
            try:
                logger.info(f"Attempting OpenAI call ({attempt+1}/{max_retries}) for '{source_filename_for_logging}' with model {model}") # MODIFIED New Log
                response = openai.ChatCompletion.create(
                    model=model,
                    messages=[
                        {"role": "system", "content": "You analyze content and extract semantic meaning."},
                        {"role": "user", "content": prompt}
                    ],
                    max_tokens=extract_length,  # Dynamic based on content
                    temperature=0.7  # Balanced between creativity and accuracy
                )
                
                # Process the raw response
                raw = response["choices"][0]["message"]["content"]
                
                # Log the raw response for debugging
                # MODIFIED: Replaced log_f with logger
                logger.debug(f"OpenAI Raw Response for '{source_filename_for_logging}' (preview): {raw[:500]}...")
                
                # Try to parse as JSON
                try:
                    parsed = json.loads(raw)
                    logger.info(f"Successfully parsed OpenAI JSON response for '{source_filename_for_logging}'") # MODIFIED New Log
                    
                    # Validate the expected fields
                    if "extract_title" not in parsed or "extract_content" not in parsed:
                        # MODIFIED: Replaced log_f with logger
                        logger.warning(f"JSON parsing successful for '{source_filename_for_logging}' but missing required fields. Got: {list(parsed.keys())}")
                        # Try a fallback approach - extract from raw text if possible
                        title_match = re.search(r'"extract_title":\s*"([^"]+)"', raw)
                        content_match = re.search(r'"extract_content":\s*"([^"]+)"', raw)
                        tags_match = re.search(r'"tags":\s*\[(.*?)\]', raw)
                        
                        title = title_match.group(1) if title_match else "Extracted Title"
                        extract = content_match.group(1) if content_match else raw
                        
                        if tags_match:
                            tags_str = tags_match.group(1)
                            tags = [tag.strip('"\'') for tag in tags_str.split(',')]
                        else:
                            tags = ["extracted"]
                            
                        return title, extract, tags
                    
                    # Get the extracted information
                    title = parsed.get("extract_title", "Untitled")
                    extract = parsed.get("extract_content", "No summary generated.")
                    tags = parsed.get("tags", ["untagged"])
                    logger.info(f"Extracted data for '{source_filename_for_logging}' - Title: '{title}', Tags: {tags}") # MODIFIED New Log
                    
                    # Basic validation
                    if not title or title == "Untitled":
                        # Try to generate a title from the first line of content
                        first_line = content.split('\n')[0].strip()
                        if len(first_line) > 5 and len(first_line) < 100:
                            title = first_line
                    
                    # Make sure extract isn't empty
                    if not extract or extract == "No summary." or extract == "No summary generated.":
                        if content_length < 1000:
                            # For short content, just use the original
                            extract = content
                        else:
                            # For longer content, use the first 500 chars
                            extract = content[:500] + "... (Extract generation failed, showing original content preview)"
                    
                    # Make sure we have some tags
                    if not tags or tags == ["untagged"]:
                        # Generate some basic tags from content
                        if "AI" in content:
                            tags.append("AI")
                        if "book" in content.lower() or "publication" in content.lower():
                            tags.append("Reading")
                        if "research" in content.lower():
                            tags.append("Research")
                        if file_type:
                            tags.append(file_type.capitalize())
                    
                    return title, extract, tags
                    
                except json.JSONDecodeError as json_err:
                    # MODIFIED: Replaced log_f with logger
                    logger.warning(f"JSON parsing error for '{source_filename_for_logging}': {str(json_err)}. Raw text (preview): {raw[:500]}...")
                    
                    # Attempt to extract meaningful content from non-JSON response
                    lines = raw.split('\n')
                    title = "Untitled"
                    for line in lines:
                        if "title" in line.lower() and ":" in line:
                            title = line.split(":", 1)[1].strip().strip('"\'')
                            break
                    
                    # Just use the raw output as the extract
                    extract = raw
                    
                    # Generate basic tags
                    tags = []
                    for line in lines:
                        if "tags" in line.lower() and ":" in line:
                            tags_part = line.split(":", 1)[1].strip()
                            tags = [t.strip().strip('",[]') for t in tags_part.split(",")]
                            break
                    
                    if not tags:
                        tags = ["extracted"]
                        if file_type:
                            tags.append(file_type.capitalize())
                    
                    return title, extract, tags
                
            except Exception as api_error:
                # MODIFIED: Replaced log_f with logger, enhanced log
                logger.error(f"OpenAI API error on attempt {attempt+1} for '{source_filename_for_logging}': {str(api_error)}")
                if attempt < max_retries - 1:
                    # MODIFIED: Removed log_f
                    time.sleep(retry_delay * (2 ** attempt))  # Exponential backoff
                else:
                    # Last attempt failed, raise the error to be caught by outer try-except
                    raise api_error
        
        # If we got here, all retries failed
        # MODIFIED: Added filename
        raise Exception(f"All OpenAI API retries failed for '{source_filename_for_logging}'")
        
    except Exception as e:
        # MODIFIED: Replaced log_f with logger, added exc_info=True
        logger.error(f"Error in get_extract for '{source_filename_for_logging}': {str(e)}", exc_info=True)
        # Removed log_f usage
        # MODIFIED: Removed print statement
        
        # Provide a more meaningful extract with the error
        error_title = "Extraction Failed"
        error_extract = f"The AI extraction process encountered an error: {str(e)}\n\nContent preview:\n{content[:300]}..."
        
        # Generate tags based on available information
        fallback_tags = ["extraction_failed"]
        if file_type:
            fallback_tags.append(file_type.capitalize())
        
        return error_title, error_extract, fallback_tags

def organize_files(input_folder="inbox", output_folder="assets", metadata_folder="metadata", api_key=None, openai_model="gpt-3.5-turbo", debug=False):
    """
    Organizes files from the input folder into the output folder.
    Extracts metadata using OpenAI and stores it in the metadata folder.
    
    Args:
        input_folder: The folder containing files to be organized.
        output_folder: The folder where organized files will be stored.
        metadata_folder: The folder where metadata files will be stored.
        api_key: OpenAI API key. If None, will use environment variable.
        openai_model: OpenAI model to use for metadata extraction.
        debug: If True, will print debug information.
    """
    if api_key:
        openai.api_key = api_key
    
    # Ensure all directories exist
    os.makedirs(input_folder, exist_ok=True)
    os.makedirs(output_folder, exist_ok=True)
    os.makedirs(metadata_folder, exist_ok=True)
    os.makedirs(os.path.join(output_folder, "sources"), exist_ok=True)
    
    # Get list of files in inbox
    inbox_files = []
    for root, _, files in os.walk(input_folder):
        for file in files:
            if file.startswith('.'):  # Skip hidden files
                continue
            inbox_files.append(os.path.join(root, file))
    
    print(f"Found {len(inbox_files)} files in inbox")
    if len(inbox_files) == 0:
        return
    
    # Process each file
    for file_index, input_path in enumerate(inbox_files):
        if debug:
            print(f"Processing {file_index+1}/{len(inbox_files)}: {input_path}")
        
        # Skip if the file no longer exists (it might have been moved by a parallel process)
        if not os.path.exists(input_path):
            continue
        
        # Skip hidden files
        if os.path.basename(input_path).startswith('.'):
            continue
        
        try:
            file_name = os.path.basename(input_path)
            file_type = infer_file_type(file_name)
            logger.info(f"Processing file: {file_name}, Inferred type: {file_type}") # DEBUG LOG

            # Define source_type_dir based on file_type or a default
            # This was missing and causing an error.
            # Assuming output_folder is the base for source types.
            source_type_dir = os.path.join(output_folder, file_type if file_type != "other" else "sources")
            os.makedirs(source_type_dir, exist_ok=True) # Ensure the directory exists

            output_path = os.path.join(source_type_dir, file_name) # Now output_path is defined correctly before potential modification

            # If a file with the same name exists, add a timestamp to make it unique
            if os.path.exists(output_path):
                name, ext = os.path.splitext(file_name)
                timestamp = int(time.time())
                file_name = f"{name}_{timestamp}{ext}"
                output_path = os.path.join(source_type_dir, file_name)
            
            # Extract text content based on file type
            text_content = None
            extraction_method = "unknown"
            is_linkedin = False

            logger.info(f"Dispatching based on file_type: '{file_type}' for file: {file_name}") # DEBUG LOG
            if file_type == "pdf":
                text_content = extract_text_from_pdf(input_path)
                extraction_method = "pdfplumber"
                if text_content: # Check if text_content is not None or empty
                    is_linkedin = "linkedin.com" in text_content.lower() or "Profile viewers" in text_content[:500]
            elif file_type == "image":
                text_content = extract_text_from_image(input_path)
                extraction_method = "ocr"
            elif file_type == "document": # This should be hit for .docx
                text_content = extract_text_from_docx(input_path)
                extraction_method = "docx"
                logger.info(f"Used extract_text_from_docx. Resulting text_content (first 50 chars): {str(text_content)[:50]}") # DEBUG LOG
            elif file_type == "presentation":
                text_content = extract_text_from_pptx(input_path)
                extraction_method = "pptx"
            elif file_type == "spreadsheet":
                text_content = extract_text_from_xlsx(input_path)
                extraction_method = "xlsx"
            elif file_type == "rtf":
                text_content = extract_text_from_rtf(input_path)
                extraction_method = "rtf"
            elif file_type == "text" and input_path.lower().endswith((".md", ".txt")): # Ensure .txt also uses this if desired, or separate
                # For .md, you might want extract_text_from_markdown if it does more than just read
                # For .txt, reading directly is fine
                if input_path.lower().endswith(".md"):
                    text_content = extract_text_from_markdown(input_path) # Assuming this exists and is preferred for .md
                    extraction_method = "markdown"
                else: # Handle .txt
                    with open(input_path, "r", encoding="utf-8", errors="ignore") as f:
                        text_content = f.read()
                    extraction_method = "textread"
            else:
                logger.warning(f"File type '{file_type}' for {file_name} not handled by specific extractors, falling back to decode.") # DEBUG LOG
                with open(input_path, "rb") as f:
                    raw_bytes = f.read()
                try:
                    text_content = raw_bytes.decode("utf-8")
                    extraction_method = "decode_utf8"
                except UnicodeDecodeError:
                    text_content = raw_bytes.decode("latin-1", errors="ignore")
                    extraction_method = "decode_latin1"
            
            logger.info(f"File: {file_name}, Extraction method: {extraction_method}") # DEBUG LOG

            # Prepare metadata
            content_preview = text_content[:2000] if text_content else ""
            
            # Generate safe filename for metadata
            metadata_filename = re.sub(r'[^\\w\\-_\\. ]', '_', file_name)
            metadata_path = os.path.join(metadata_folder, f"{os.path.splitext(metadata_filename)[0]}.md")
            
            # No longer need response_type here as get_extract handles different content types
            # response_type = "extract"
            # if is_linkedin:
            #     response_type = "linkedin_post"
            
            # Generate metadata using OpenAI via get_extract
            today = time.strftime("%Y-%m-%d")

            if text_content and len(text_content.strip()) > 0:
                logger.info(f"Attempting AI enrichment (get_extract) for {file_name}")
                try:
                    # get_extract returns: title_from_ai, summary_from_ai, tags_from_ai
                    returned_title, returned_summary, returned_tags = get_extract(
                        text_content,
                        file_type=file_type,
                        urls_metadata=None, # TODO: Consider populating urls_metadata if enrich_urls is called prior
                        is_linkedin=is_linkedin,
                        source_filename_for_logging=file_name
                    )

                    # Check if get_extract provided a valid, non-error summary
                    if returned_summary and \
                       "Extract failed" not in returned_summary and \
                       "Missing API Key" not in returned_summary and \
                       "No summary" not in returned_summary and \
                       len(returned_summary.strip()) > 10: # Basic check for meaningful summary
                        
                        logger.info(f"Successfully received AI-generated content from get_extract for {file_name}")
                        
                        metadata_dict = {
                            "title": returned_title if returned_title else f"AI Processed: {file_name}",
                            "author": "Unknown",
                            "date": today,
                            "category": "Other", # Can be refined by AI or later review
                            "tags": returned_tags if returned_tags else ["processed"],
                            "extract_content": returned_summary,  # AI-generated summary
                            "parse_status": "enriched",
                            "extraction_method": extraction_method,
                            "file_type": file_type,
                            "source": file_name,
                            "reviewed": False,
                            "reprocess_status": "none",
                            "reprocess_rounds": "0",
                            "source_url": None
                        }
                        
                        if not isinstance(metadata_dict.get("tags"), list):
                            tags_val = metadata_dict.get("tags")
                            metadata_dict["tags"] = [str(tags_val)] if tags_val else []

                        if is_linkedin:
                            if "linkedin" not in metadata_dict["tags"]:
                                metadata_dict["tags"].append("linkedin")
                            if "social_media" not in metadata_dict["tags"]:
                                metadata_dict["tags"].append("social_media")
                            metadata_dict["category"] = "Social Media"

                        current_text_content_for_file = text_content if text_content is not None else ""
                        post = frontmatter.Post(current_text_content_for_file, **metadata_dict)
                        
                        with open(metadata_path, 'wb') as f:
                            frontmatter.dump(post, f)
                            
                        shutil.move(input_path, output_path)
                        logger.info(f"Organized with AI summary: {file_name} -> {output_path}")

                    else: 
                        logger.warning(f"get_extract for '{file_name}' did not return a usable summary or failed. Summary: '{returned_summary}'. Proceeding with blank extract.")
                        metadata_dict = {
                            "title": f"Needs Review: {file_name}",
                            "author": "Unknown",
                            "date": today,
                            "category": "Other",
                            "tags": ["needs_review", "enrichment_failed"],
                            "extract_content": "",  # Blank as per user request
                            "parse_status": "enrichment_failed",
                            "extraction_method": extraction_method,
                            "file_type": file_type,
                            "source": file_name,
                            "reviewed": False,
                            "reprocess_status": "none",
                            "reprocess_rounds": "0",
                            "source_url": None
                        }
                        current_text_content_for_file = text_content if text_content is not None else ""
                        post = frontmatter.Post(current_text_content_for_file, **metadata_dict)
                        with open(metadata_path, 'wb') as f:
                            frontmatter.dump(post, f)
                        shutil.move(input_path, output_path)
                        logger.info(f"Organized '{file_name}' with enrichment_failed status and blank extract.")
                        
                except Exception as ge_exc: 
                    logger.error(f"Error during get_extract call for '{file_name}': {str(ge_exc)}. Proceeding with blank extract.", exc_info=True)
                    metadata_dict = {
                        "title": f"Error Processing: {file_name}",
                        "author": "Unknown",
                        "date": today,
                        "category": "Other",
                        "tags": ["error", "enrichment_failed"],
                        "extract_content": "",  # Blank as per user request
                        "parse_status": "enrichment_failed",
                        "extraction_method": extraction_method,
                        "file_type": file_type,
                        "source": file_name,
                        "reviewed": False,
                        "reprocess_status": "needs_review", # Flag for potential reprocessing
                        "reprocess_rounds": "0",
                        "source_url": None
                    }
                    current_text_content_for_file = text_content if text_content is not None else ""
                    post = frontmatter.Post(current_text_content_for_file, **metadata_dict)
                    with open(metadata_path, 'wb') as f:
                        frontmatter.dump(post, f)
                    shutil.move(input_path, output_path)
                    logger.info(f"Organized '{file_name}' with error during enrichment and blank extract.")
            else: 
                logger.info(f"No content extracted from \'{file_name}\'. Using basic_metadata.")
                basic_metadata(input_path, output_path, metadata_path, file_name, file_type, text_content, extraction_method)
                
        except Exception as e:
            logger.error(f"Error processing {input_path} in organize_files: {str(e)}", exc_info=True) # Add exc_info for traceback            # Move to error folder
            error_folder = os.path.join(output_folder, "errors")
            os.makedirs(error_folder, exist_ok=True)
            error_path = os.path.join(error_folder, os.path.basename(input_path))
            try:
                shutil.move(input_path, error_path)
            except:
                # If moving fails, just delete it
                os.remove(input_path)
    
    print("File organization complete")
    
    # Log out the final status for debugging
    processed_count = len(inbox_files)
    errors_count = len([f for f in os.listdir(output_folder) if f.endswith(".error")]) if os.path.exists(output_folder) else 0
    metadata_count = len([f for f in os.listdir(metadata_folder) if f.endswith(".md")]) if os.path.exists(metadata_folder) else 0
    
    logger.info(f"Organize files summary - Input folder: {input_folder}, Files processed: {processed_count}, "
                f"Metadata files created: {metadata_count}, Errors: {errors_count}")
    
    # Return success and failure metrics
    return {
        "success_count": metadata_count,
        "processed_count": processed_count,
        "failed_files": [(os.path.basename(f), "Processing failed") for f in os.listdir(output_folder) if f.endswith(".error")] if os.path.exists(output_folder) else []
    }
    
def basic_metadata(input_path, output_path, metadata_path, file_name, file_type, text_content, extraction_method):
    logger.warning(f"Executing basic_metadata for file '{file_name}' due to previous error or lack of content. Extraction method: {extraction_method}") # New Log
    today = time.strftime("%Y-%m-%d")
    
    # Ensure text_content is a string for safe operations
    current_text_content = text_content if text_content is not None else ""
    
    # Create basic metadata
    metadata_dict = {
        "title": f"Unprocessed: {file_name}",
        "author": "Unknown",
        "date": today,
        "category": "Other",
        "tags": ["unprocessed"],
        # Use current_text_content which is guaranteed to be a string
        "extract_content": current_text_content[:500] + "..." if len(current_text_content) > 500 else current_text_content,
        "parse_status": "basic",
        "extraction_method": extraction_method,
        "file_type": file_type,
        "source": file_name,
        "reviewed": False,
        "reprocess_status": "none",
        "reprocess_rounds": "0",
        "source_url": None
    }
    
    # Create frontmatter post
    post = frontmatter.Post(current_text_content, **metadata_dict) # Use current_text_content
    
    # Save metadata file
    with open(metadata_path, 'wb') as f:
        frontmatter.dump(post, f)
        
    # Move source file to organized location
    shutil.move(input_path, output_path)